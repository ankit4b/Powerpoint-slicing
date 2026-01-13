"""
PowerPoint Slicer - Split a PowerPoint file into individual slides
Each slide will be saved as a separate PowerPoint file and as images.
Cross-platform solution compatible with Databricks and Linux environments.
"""

import os
import subprocess
import shutil
from pathlib import Path
from pptx import Presentation
from PIL import Image
import platform


def export_slides_as_images(input_file, output_dir=None, image_format='png'):
    """
    Export each slide as an image file.
    Falls back to Windows COM if LibreOffice not available.
    
    Args:
        input_file (str): Path to the input PowerPoint file
        output_dir (str, optional): Directory to save images. 
                                   Defaults to 'output/images' folder.
        image_format (str): Image format - 'png', 'jpg', or 'jpeg'
    
    Returns:
        list: List of paths to the created image files
    """
    # Set up output directory
    if output_dir is None:
        output_dir = os.path.join("output", "images")
    
    Path(output_dir).mkdir(parents=True, exist_ok=True)
    
    # Convert to absolute path
    input_file = os.path.abspath(input_file)
    output_dir = os.path.abspath(output_dir)
    
    print(f"Exporting slides as images from: {input_file}")
    
    # Check if we're on Windows and LibreOffice is not available
    is_windows = platform.system() == 'Windows'
    
    # Try to find LibreOffice
    libreoffice_paths = [
        'libreoffice',  # Linux/Mac
        'soffice',      # Alternative name
        '/usr/bin/libreoffice',
        '/usr/bin/soffice',
        'C:\\Program Files\\LibreOffice\\program\\soffice.exe',
        'C:\\Program Files (x86)\\LibreOffice\\program\\soffice.exe',
    ]
    
    libreoffice_cmd = None
    for path in libreoffice_paths:
        if shutil.which(path) or os.path.exists(path):
            libreoffice_cmd = path
            break
    
    # If on Windows and no LibreOffice, use Windows COM
    if is_windows and not libreoffice_cmd:
        print("LibreOffice not found. Using Windows PowerPoint COM automation...")
        return _export_slides_windows_com(input_file, output_dir, image_format)
    
    # Use LibreOffice method
    if not libreoffice_cmd:
        raise RuntimeError(
            "LibreOffice not found. Please install LibreOffice:\n"
            "- Ubuntu/Debian: sudo apt-get install libreoffice\n"
            "- CentOS/RHEL: sudo yum install libreoffice\n"
            "- macOS: brew install --cask libreoffice\n"
            "- Windows: Download from https://www.libreoffice.org/download/"
        )
    
    return _export_slides_libreoffice(input_file, output_dir, image_format, libreoffice_cmd)


def _export_slides_windows_com(input_file, output_dir, image_format='png'):
    """Export slides using Windows PowerPoint COM automation."""
    try:
        import win32com.client
    except ImportError:
        raise ImportError(
            "pywin32 not installed. Install it with: pip install pywin32\n"
            "Alternatively, install LibreOffice for cross-platform support."
        )
    
    # Load presentation to get slide count
    prs = Presentation(input_file)
    total_slides = len(prs.slides)
    input_name = Path(input_file).stem
    
    # Initialize PowerPoint
    powerpoint = win32com.client.Dispatch("Powerpoint.Application")
    powerpoint.Visible = 1
    
    try:
        # Open the presentation
        presentation = powerpoint.Presentations.Open(input_file, WithWindow=False)
        
        created_images = []
        img_ext = image_format.lower()
        if img_ext == 'jpeg':
            img_ext = 'jpg'
        
        # Export each slide
        for slide_idx in range(1, total_slides + 1):
            slide = presentation.Slides(slide_idx)
            image_path = os.path.join(output_dir, f"{input_name}_slide_{slide_idx}.{img_ext}")
            
            # Export slide as image
            if img_ext == 'png':
                slide.Export(image_path, "PNG")
            elif img_ext == 'jpg':
                slide.Export(image_path, "JPG")
            
            created_images.append(image_path)
            print(f"Created image: {image_path}")
        
        # Close presentation
        presentation.Close()
        print(f"\n✓ Successfully created {len(created_images)} image files in '{output_dir}'")
        
        return created_images
        
    finally:
        # Quit PowerPoint
        powerpoint.Quit()


def _export_slides_libreoffice(input_file, output_dir, image_format, libreoffice_cmd):
    """Export slides using LibreOffice (cross-platform)."""
    # Load presentation to get slide count
    prs = Presentation(input_file)
    total_slides = len(prs.slides)
    input_name = Path(input_file).stem
    
def _export_slides_libreoffice(input_file, output_dir, image_format, libreoffice_cmd):
    """Export slides using LibreOffice (cross-platform)."""
    # Load presentation to get slide count
    prs = Presentation(input_file)
    total_slides = len(prs.slides)
    input_name = Path(input_file).stem
    
    # Create a temporary directory for PDF conversion
    temp_dir = os.path.join(output_dir, "temp_pdf")
    Path(temp_dir).mkdir(parents=True, exist_ok=True)
    
    try:
        # Step 1: Convert PPTX to PDF using LibreOffice
        pdf_output = os.path.join(temp_dir, f"{input_name}.pdf")
        
        print(f"Using LibreOffice: {libreoffice_cmd}")
        
        # Convert PPTX to PDF
        print("Converting PPTX to PDF...")
        cmd = [
            libreoffice_cmd,
            '--headless',
            '--convert-to', 'pdf',
            '--outdir', temp_dir,
            input_file
        ]
        
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=300)
        
        if result.returncode != 0:
            raise RuntimeError(f"LibreOffice conversion failed: {result.stderr}")
        
        if not os.path.exists(pdf_output):
            raise RuntimeError(f"PDF file not created: {pdf_output}")
        
        # Step 2: Convert PDF to images using pdf2image
        print("Converting PDF pages to images...")
        try:
            from pdf2image import convert_from_path
            
            # Convert PDF to images
            images = convert_from_path(pdf_output, dpi=300)
            
            created_images = []
            img_ext = image_format.lower()
            if img_ext == 'jpeg':
                img_ext = 'jpg'
            
            for idx, image in enumerate(images, start=1):
                image_path = os.path.join(output_dir, f"{input_name}_slide_{idx}.{img_ext}")
                
                # Save image
                if img_ext == 'jpg':
                    image.save(image_path, 'JPEG', quality=95)
                else:
                    image.save(image_path, 'PNG')
                
                created_images.append(image_path)
                print(f"Created image: {image_path}")
            
            print(f"\n✓ Successfully created {len(created_images)} image files in '{output_dir}'")
            return created_images
            
        except ImportError:
            raise ImportError(
                "pdf2image not installed. Install it with: pip install pdf2image\n"
                "Also install poppler:\n"
                "- Ubuntu/Debian: sudo apt-get install poppler-utils\n"
                "- CentOS/RHEL: sudo yum install poppler-utils\n"
                "- macOS: brew install poppler\n"
                "- Windows: Download from https://github.com/oschwartz10612/poppler-windows/releases/"
            )
    
    finally:
        # Clean up temporary directory
        if os.path.exists(temp_dir):
            shutil.rmtree(temp_dir)


def split_pptx(input_file, output_dir=None):
    """
    Split a PowerPoint file into individual slides.
    
    Args:
        input_file (str): Path to the input PowerPoint file
        output_dir (str, optional): Directory to save output files. 
                                   Defaults to 'output/pptx' folder.
    
    Returns:
        list: List of paths to the created PowerPoint files
    """
    # Validate input file
    if not os.path.exists(input_file):
        raise FileNotFoundError(f"Input file not found: {input_file}")
    
    # Set up output directory
    if output_dir is None:
        output_dir = os.path.join("output", "pptx")
    
    Path(output_dir).mkdir(parents=True, exist_ok=True)
    
    # Load the presentation to get slide count
    print(f"Loading presentation: {input_file}")
    prs = Presentation(input_file)
    total_slides = len(prs.slides)
    print(f"Total slides found: {total_slides}")
    
    created_files = []
    input_name = Path(input_file).stem
    
    # Process each slide
    for slide_idx in range(total_slides):
        # Create output filename
        output_file = os.path.join(output_dir, f"{input_name}_slide_{slide_idx + 1}.pptx")
        
        # Create a fresh copy of the presentation for each slide
        # Load the presentation fresh each time
        temp_prs = Presentation(input_file)
        
        # Create list of slide indices to remove (all except the current one)
        slides_to_remove = list(range(total_slides))
        slides_to_remove.remove(slide_idx)
        
        # Remove slides in reverse order to maintain indices
        for idx in reversed(slides_to_remove):
            rId = temp_prs.slides._sldIdLst[idx].rId
            temp_prs.part.drop_rel(rId)
            del temp_prs.slides._sldIdLst[idx]
        
        # Save the presentation with only one slide
        temp_prs.save(output_file)
        created_files.append(output_file)
        print(f"Created: {output_file}")
    
    print(f"\n✓ Successfully created {total_slides} PowerPoint files in '{output_dir}'")
    return created_files


def split_pptx_with_images(input_file, output_pptx_dir=None, output_images_dir=None, image_format='png'):
    """
    Split PowerPoint file into individual slides AND export as images.
    
    Args:
        input_file (str): Path to the input PowerPoint file
        output_pptx_dir (str, optional): Directory to save PPTX files
        output_images_dir (str, optional): Directory to save image files
        image_format (str): Image format - 'png', 'jpg', or 'jpeg'
    
    Returns:
        dict: Dictionary with 'pptx_files' and 'image_files' lists
    """
    print("=" * 60)
    print("PowerPoint Slicer - Generating PPTs and Images")
    print("=" * 60)
    print()
    
    # Generate individual PPTX files
    print("Step 1: Creating individual PPTX files...")
    pptx_files = split_pptx(input_file, output_pptx_dir)
    print()
    
    # Generate images
    print("Step 2: Creating slide images...")
    image_files = export_slides_as_images(input_file, output_images_dir, image_format)
    print()
    
    print("=" * 60)
    print(f"✓ Complete! Generated {len(pptx_files)} PPTX files and {len(image_files)} images")
    print("=" * 60)
    
    return {
        'pptx_files': pptx_files,
        'image_files': image_files,
        'total_slides': len(pptx_files)
    }


def main():
    """Main function to run the PowerPoint slicer."""
    import argparse
    
    parser = argparse.ArgumentParser(
        description="Split a PowerPoint file into individual slides and images"
    )
    parser.add_argument(
        "input_file",
        help="Path to the input PowerPoint file (.pptx)"
    )
    parser.add_argument(
        "-o", "--output",
        default="output",
        help="Base output directory (default: 'output')"
    )
    parser.add_argument(
        "--with-images",
        action="store_true",
        help="Also generate images for each slide"
    )
    parser.add_argument(
        "--image-format",
        default="png",
        choices=["png", "jpg", "jpeg"],
        help="Image format (default: png)"
    )
    
    args = parser.parse_args()
    
    try:
        if args.with_images:
            # Generate both PPTX and images
            pptx_dir = os.path.join(args.output, "pptx")
            images_dir = os.path.join(args.output, "images")
            split_pptx_with_images(args.input_file, pptx_dir, images_dir, args.image_format)
        else:
            # Generate only PPTX files
            split_pptx(args.input_file, args.output)
    except Exception as e:
        print(f"Error: {e}")
        return 1
    
    return 0


if __name__ == "__main__":
    exit(main())
